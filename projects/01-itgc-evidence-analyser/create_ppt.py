#!/usr/bin/env python3
"""Generate the ITGC Evidence Analyser manager presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Vodafone brand colours
VODAFONE_RED = RGBColor(0xE6, 0x00, 0x00)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY = RGBColor(0x6E, 0x76, 0x87)
LIGHT_GREY = RGBColor(0xE8, 0xEA, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x5B, 0x8D, 0xEF)
PASS_GREEN = RGBColor(0x26, 0xC9, 0x63)
BG_DARK = RGBColor(0x05, 0x07, 0x0A)
SURFACE = RGBColor(0x0A, 0x0E, 0x14)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=0, top=0, width=None, height=Inches(0.06), color=VODAFONE_RED):
    if width is None:
        width = prs.slide_width
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        p.bullet = True
    return tf


def add_card(slide, left, top, width, height, title, body, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = RGBColor(0x1A, 0x23, 0x32)
    shape.line.width = Pt(1)

    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), Inches(0.4),
                 title, font_size=14, color=title_color, bold=True)
    # Body
    add_text_box(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.5), height - Inches(0.7),
                 body, font_size=11, color=MID_GREY)
    return shape


def add_stat_card(slide, left, top, width, height, number, label, accent=VODAFONE_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = RGBColor(0x1A, 0x23, 0x32)
    shape.line.width = Pt(1)

    add_accent_bar(slide, left, top, width, Inches(0.04), accent)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.6),
                 number, font_size=28, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.75), width - Inches(0.4), Inches(0.3),
                 label, font_size=10, color=MID_GREY)
    return shape


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(2.8), Inches(2), Inches(0.06))
add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.2),
             "ITGC Evidence Analyser", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.0), Inches(10), Inches(0.6),
             "AI-Powered Audit Evidence Assessment for Vodafone ITGC Controls", font_size=20, color=MID_GREY)
add_text_box(slide, Inches(1), Inches(3.7), Inches(10), Inches(0.5),
             "GRC Engineering Portfolio  |  Seth Odoi Asare  |  April 2026", font_size=14, color=MID_GREY)
add_accent_bar(slide, Inches(1), Inches(4.3), Inches(4), Inches(0.02), MID_GREY)

# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(0.8), Inches(2), Inches(0.05))
add_text_box(slide, Inches(1), Inches(0.3), Inches(4), Inches(0.6),
             "The Problem", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "Traditional ITGC assurance is manual, slow, and inconsistent across 32 markets.", font_size=16, color=MID_GREY)

problems = [
    "58 controls × 32 markets = 1,856 possible assessments — impossible to sample thoroughly",
    "Auditors spend 60-70% of time collecting and formatting evidence, not analysing it",
    "Inconsistent verdicts across markets — same evidence, different auditor, different outcome",
    "No centralized trail of what was assessed, when, by whom, and against which evidence",
    "Junior auditors lack the experience to draft audit-grade findings and risk ratings",
    "Evidence re-examination requires re-reading all files — no conversational interface",
]
add_bullet_list(slide, Inches(1), Inches(2.0), Inches(7), Inches(4.5), problems, font_size=15, color=LIGHT_GREY, spacing=Pt(12))

add_stat_card(slide, Inches(9), Inches(1.5), Inches(3.2), Inches(1.2), "1,856", "Possible Assessments", VODAFONE_RED)
add_stat_card(slide, Inches(9), Inches(3.0), Inches(3.2), Inches(1.2), "~70%", "Time on Evidence Prep", VODAFONE_RED)
add_stat_card(slide, Inches(9), Inches(4.5), Inches(3.2), Inches(1.2), "32", "Markets Needing Coverage", VODAFONE_RED)

# ============================================================
# SLIDE 3 — THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(0.8), Inches(2), Inches(0.05))
add_text_box(slide, Inches(1), Inches(0.3), Inches(4), Inches(0.6),
             "The Solution", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "An AI-powered platform that automates evidence assessment with audit-grade precision.", font_size=16, color=MID_GREY)

solutions = [
    "Upload evidence once — the AI assessor evaluates it against the exact D/E statements for any control",
    "Consistent, repeatable verdicts across all 32 markets using the same assessment methodology",
    "Automated draft findings in Vodafone audit format — title, observation, criteria, risk impact, recommendation",
    "Full audit trail: who assessed what, when, which market, which samples, every evidence file logged",
    "AI chatbot that can re-read evidence and amend verdicts if the auditor catches something missed",
    "Multi-user workspace isolation — each auditor sees only their own assessments, login audit log",
]
add_bullet_list(slide, Inches(1), Inches(2.0), Inches(7.5), Inches(4.5), solutions, font_size=14, color=LIGHT_GREY, spacing=Pt(10))

add_stat_card(slide, Inches(9.5), Inches(1.5), Inches(3), Inches(1.2), "58", "Controls Covered", PASS_GREEN)
add_stat_card(slide, Inches(9.5), Inches(3.0), Inches(3), Inches(1.2), "10", "ITGC Domains", ACCENT_BLUE)
add_stat_card(slide, Inches(9.5), Inches(4.5), Inches(3), Inches(1.2), "5 min", "Per Assessment", PASS_GREEN)

# ============================================================
# SLIDE 4 — KEY FEATURES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(0.8), Inches(2), Inches(0.05))
add_text_box(slide, Inches(1), Inches(0.3), Inches(4), Inches(0.6),
             "Key Features", font_size=32, color=WHITE, bold=True)

features = [
    ("AI-Powered Assessment", "Claude Sonnet 4.6 evaluates evidence against exact D and E statements. Returns structured verdicts (PASS/PARTIAL/FAIL), confidence scores, risk ratings, gap analysis, and draft audit findings."),
    ("Multi-Market Architecture", "32 Vodafone subsidiaries pre-loaded. Each assessment scoped to a specific market with samples in scope (e.g., Phones, Tablets for DRC). Samples persist across assessments for repeatable testing."),
    ("Evidence Intelligence", "Upload PDFs, DOCX, XLSX, CSV, screenshots, or paste text. The assessor extracts, catalogues, and cross-references every piece of evidence against control requirements."),
    ("AI Audit Chatbot", "Conversational interface tied to assessments. Ask the assessor to re-read evidence, explain verdicts, or amend findings if you spot something it missed. Full tool-use integration."),
    ("Audit Trail & Workspaces", "Every assessment logged with user, timestamp, market, samples, and evidence files. Multi-user with JWT authentication. Each auditor gets their own isolated workspace."),
    ("Professional PDF Export", "Vodafone-branded audit reports with market context, samples in scope, full finding narratives, evidence inventory tables, and remediation guidance — ready for control owner sign-off."),
]

for i, (title, body) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(1 + col * 3.8)
    top = Inches(1.5 + row * 2.8)
    add_card(slide, left, top, Inches(3.5), Inches(2.5), title, body)

# ============================================================
# SLIDE 5 — INCORPORATION INTO TRADITIONAL ASSURANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(0.8), Inches(2), Inches(0.05))
add_text_box(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.6),
             "How It Fits Into Traditional Assurance", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "The tool does not replace auditors — it amplifies them. Here is the new workflow:", font_size=16, color=MID_GREY)

steps = [
    ("1. Plan", "Auditor selects Market + Control + Samples in Scope (e.g., Czech Republic → Mobile Device Policy → Phones, Tablets)"),
    ("2. Collect", "Evidence owner uploads policy docs, screenshots, MDM reports, configuration exports"),
    ("3. AI Assesses", "Claude evaluates evidence against every D/E statement. Produces verdict, gaps, risk rating, and draft finding in under 60 seconds"),
    ("4. Auditor Reviews", "Auditor examines the AI output. If something was missed, uses the chatbot to re-read evidence and amend the verdict"),
    ("5. Approve & Export", "Final verdict confirmed. One-click PDF export — professional audit report ready for control owner"),
    ("6. Track", "All assessments logged in the audit trail. Filter by market, control, or verdict. Evidence freshness tracked across all markets"),
]

for i, (step_title, step_body) in enumerate(steps):
    left = Inches(1 + (i % 3) * 3.8)
    top = Inches(2.0 + (i // 3) * 2.6)
    # Step number
    add_text_box(slide, left, top, Inches(0.5), Inches(0.4),
                 step_title.split(".")[0], font_size=24, color=VODAFONE_RED, bold=True)
    # Title
    add_text_box(slide, left + Inches(0.55), top, Inches(3), Inches(0.4),
                 step_title.split(" ", 1)[1], font_size=15, color=WHITE, bold=True)
    # Body
    add_text_box(slide, left + Inches(0.55), top + Inches(0.4), Inches(3.2), Inches(1.6),
                 step_body, font_size=11, color=MID_GREY)

# ============================================================
# SLIDE 6 — BENEFITS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(0.8), Inches(2), Inches(0.05))
add_text_box(slide, Inches(1), Inches(0.3), Inches(4), Inches(0.6),
             "Benefits", font_size=32, color=WHITE, bold=True)

benefits = [
    ("80% Faster", "Assessment time reduced from hours to minutes. One auditor can cover all 32 markets.", PASS_GREEN),
    ("Consistent", "Same AI methodology applied to every assessment. No auditor-to-auditor variance in verdicts.", ACCENT_BLUE),
    ("Scalable", "Add markets, controls, or samples without hiring. The AI assessor scales with your portfolio.", ACCENT_BLUE),
    ("Audit-Ready", "Every output includes a full evidence inventory, requirement assessment table, and draft finding in Vodafone format.", PASS_GREEN),
    ("Traceable", "Complete audit trail — who assessed what, when, which evidence. Non-repudiation via user accounts.", ACCENT_BLUE),
    ("Cost-Effective", "Uses Claude API with prompt caching (~90% input cost reduction). No per-seat licence fees. Open-source.", PASS_GREEN),
]

for i, (title, body, accent) in enumerate(benefits):
    col = i % 3
    row = i // 3
    left = Inches(1 + col * 3.8)
    top = Inches(1.5 + row * 2.8)
    add_card(slide, left, top, Inches(3.5), Inches(2.5), title, body, title_color=accent)

# ============================================================
# SLIDE 7 — ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(0.8), Inches(2), Inches(0.05))
add_text_box(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.6),
             "Technical Architecture", font_size=32, color=WHITE, bold=True)

arch_text = (
    "FastAPI Backend (Python)  •  Next.js 16 Frontend (React 19)  •  Anthropic Claude API (Sonnet 4.6)\n"
    "SQLite Database  •  JWT Authentication  •  Docker Container  •  Azure Container Apps\n\n"
    "The frontend communicates with the backend via REST API. Evidence files are uploaded,\n"
    "extracted (PDF, DOCX, XLSX, CSV, images), and evaluated by Claude against the exact\n"
    "D/E statements for the selected control. Results are persisted to SQLite with full\n"
    "audit trail metadata (user, market, samples, timestamp, evidence references).\n\n"
    "The chatbot uses Claude's function-calling capability with 5 tools:\n"
    "list_assessments, get_assessment, reread_evidence, modify_verdict, get_control_detail\n\n"
    "All 52 automated tests pass. Docker multi-stage build. Single-container deployment\n"
    "with Nginx reverse proxy and Supervisor process management."
)
add_text_box(slide, Inches(1), Inches(1.5), Inches(7.5), Inches(5.5),
             arch_text, font_size=13, color=LIGHT_GREY)

add_stat_card(slide, Inches(9.5), Inches(1.5), Inches(3), Inches(1.0), "52", "Tests Passing", PASS_GREEN)
add_stat_card(slide, Inches(9.5), Inches(2.8), Inches(3), Inches(1.0), "7", "DB Tables", ACCENT_BLUE)
add_stat_card(slide, Inches(9.5), Inches(4.1), Inches(3), Inches(1.0), "20+", "API Endpoints", ACCENT_BLUE)
add_stat_card(slide, Inches(9.5), Inches(5.4), Inches(3), Inches(1.0), "1", "Docker Image", PASS_GREEN)

# ============================================================
# SLIDE 8 — ROADMAP & NEXT STEPS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(1), Inches(0.8), Inches(2), Inches(0.05))
add_text_box(slide, Inches(1), Inches(0.3), Inches(4), Inches(0.6),
             "Next Steps", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "Immediate deployment and future roadmap for the GRC portfolio.", font_size=16, color=MID_GREY)

roadmap = [
    "This week: Deploy ITGC Evidence Analyser to Azure — available to all 32 markets",
    "Next sprint: Evidence Collection Automation — scripts to auto-collect logs, configs, screenshots",
    "In pipeline: Control Coverage Mapper, Risk Register, Cloud Posture Snapshot, Vendor Scorer",
    "Target: Full GRC automation platform covering the entire assurance lifecycle by Q3 2026",
]
add_bullet_list(slide, Inches(1), Inches(2.0), Inches(7.5), Inches(3.5), roadmap, font_size=14, color=LIGHT_GREY, spacing=Pt(14))

# Call to action
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.5), Inches(5), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = VODAFONE_RED
shape.line.fill.background()
add_text_box(slide, Inches(1.3), Inches(5.65), Inches(4.4), Inches(0.9),
             "Ready to deploy. Request a demo.\nContact: Seth Odoi Asare — Cyber Assurance Manager", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(OUTPUT_DIR, "ITGC_Evidence_Analyser_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"File size: {os.path.getsize(output_path)} bytes")
